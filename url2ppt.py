# -*- coding: utf-8 -*-
import os
import sys
import time
import logging
import json
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN  
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image, UnidentifiedImageError
import re
import pytesseract
from typing import Optional, List, Dict, Tuple, Any
from urllib.parse import urljoin
from newspaper import Article, ArticleException
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import backoff

# 配置日志系统
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('web_to_ppt.log'),
        logging.StreamHandler(sys.stdout)
    ]
)

# 全局常量
MIN_IMAGE_WIDTH = 300  # ppt中图片最小尺寸要求
MIN_IMAGE_HEIGHT = 250
MAX_IMAGES_PER_PPT = 6  # ppt中的图片数量限制
MAX_RETRIES = 3
API_TIMEOUT = 60  # Deepseek API调用超时时间
IMAGE_DOWNLOAD_TIMEOUT = 30
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"
MAX_CONTENT_LENGTH = 12000  # 增加最大内容长度
MIN_SECTION_LENGTH = 300  # 最小段落长度

class PPTGenerator:
    def __init__(self):
        self.api_key = self.get_api_key()
        self.headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8'
        }
        
    def get_api_key(self) -> str:
        """获取并验证Deepseek API密钥"""
        api_key = os.getenv('DEEPSEEK_API_KEY', '')
        if not api_key or not api_key.startswith('sk-'):
            raise ValueError("无效的DeepSeek API密钥。请设置 DEEPSEEK_API_KEY 环境变量")
        return api_key

    @backoff.on_exception(backoff.expo, 
                         (requests.exceptions.RequestException, 
                          requests.exceptions.Timeout),
                         max_tries=MAX_RETRIES)
    def call_deepseek_api(self, system_prompt: str, user_prompt: str, 
                         max_tokens: int = 4000, temperature: float = 0.3) -> Optional[str]:
        """调用Deepseek API，对文章内容主题和要点进行总结"""
        headers = {
            "Authorization": f"Bearer {self.api_key}", 
            "Content-Type": "application/json"
        }
        
        data = {
            "model": "deepseek-chat",
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            "max_tokens": max_tokens,
            "temperature": temperature
        }
        
        try:
            response = requests.post(
                DEEPSEEK_API_URL, 
                headers=headers, 
                json=data, 
                timeout=API_TIMEOUT
            )
            response.raise_for_status()
            
            result = response.json()
            if 'choices' in result and result['choices']:
                return result['choices'][0]['message']['content'].strip()
            return None
        except Exception as e:
            logging.error(f"API调用失败: {str(e)}")
            raise

    def download_image(self, img_url: str, base_url: str) -> Optional[Tuple[BytesIO, Tuple[int, int]]]:
        """图片下载器，过滤SVG格式，并支持webp格式转换"""
        if not img_url or img_url.startswith('data:image'):
            return None

        # 如果 URL 中包含 svg 字样，则跳过
        if 'svg' in img_url.lower():
            logging.info(f"跳过SVG图片: {img_url}")
            return None

        try:
            absolute_url = urljoin(base_url, img_url.strip())
            response = requests.get(
                absolute_url, 
                headers=self.headers, 
                timeout=IMAGE_DOWNLOAD_TIMEOUT, 
                stream=True
            )
            response.raise_for_status()
            
            # 验证图片类型
            content_type = response.headers.get('Content-Type', '')
            # 如果Content-Type提示为svg或非图片，则跳过
            if not content_type or not content_type.startswith('image/') or 'svg' in content_type.lower():
                logging.info(f"跳过非图片或SVG图片，Content-Type={content_type}: {absolute_url}")
                return None

            # 下载图片数据
            img_data = BytesIO()
            for chunk in response.iter_content(chunk_size=8192):
                img_data.write(chunk)
            img_data.seek(0)
            
            # 针对 webp 格式处理：如果Content-Type中包含webp或者URL后缀为webp，则转换为PNG
            if 'webp' in content_type.lower() or '.webp' in absolute_url.lower() or '.awebp' in absolute_url.lower():
                try:
                    with Image.open(img_data) as img:
                        img = img.convert('RGB')
                        converted_data = BytesIO()
                        img.save(converted_data, format='PNG')
                        converted_data.seek(0)
                        width, height = img.size
                        if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                            return None
                        return converted_data, (width, height)
                except Exception as e:
                    logging.error(f"处理webp图片失败 {absolute_url}: {str(e)}")
                    return None
            
            # 对其他格式进行处理：尝试打开图片验证是否有效
            try:
                with Image.open(img_data) as img:
                    img.verify()  # 验证图片完整性
                    img_data.seek(0)
                    width, height = img.size
                    if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                        return None
                    return img_data, (width, height)
            except Exception as e:
                logging.error(f"图片打开失败 {absolute_url}: {str(e)}")
                return None
                
        except Exception as e:
            logging.error(f"图片下载失败 {img_url}: {str(e)}")
            return None


    def parse_article(self, url: str) -> Dict[str, Any]:
        """网页解析器，优先使用newspaper3k，再使用BeautifulSoup进一步提取正文和图片"""
        try:
            # 优先使用newspaper3k
            article = Article(url, request_headers=self.headers)
            article.download()
            article.parse()
            
            text = self.clean_text(article.text)
            images = list(article.images) if article.images else []
            
            # 如果使用newspaper3k提取的内容不够完整，则采用BeautifulSoup补充
            if not text or len(text) < 500:
                raise ArticleException("内容不足")
                
            # 使用BeautifulSoup进一步提取正文区域内的图片
            response = requests.get(url, headers=self.headers, timeout=20)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'lxml')
            main_container = self.extract_main_content(soup)
            # 提取正文图片（从 <img> 标签、figure标签等提取）
            content_images = []
            if main_container:
                for img in main_container.find_all('img', src=True):
                    src = img.get('src')
                    if src:
                        absolute = urljoin(url, src.strip())
                        if absolute not in images and absolute not in content_images:
                            content_images.append(absolute)
            # 合并两者，保证图片相关性
            images = list(set(images + content_images))
            
            return {
                'title': article.title.strip() if article.title else (soup.title.string.strip() if soup.title else "无标题"),
                'text': text,
                'html': response.text,
                'top_image': images[0] if images else None,
                'images': images
            }
            
        except Exception as e:
            logging.error(f"newspaper3k解析失败，尝试BeautifulSoup方案: {str(e)}")
            # 备用方案：BeautifulSoup解析
            try:
                response = requests.get(url, headers=self.headers, timeout=20)
                response.raise_for_status()
                soup = BeautifulSoup(response.content, 'lxml')
                
                main_container = self.extract_main_content(soup)
                text = self.clean_text(main_container.get_text(separator=" ", strip=True)) if main_container else ""
                images = []
                # 尽量提取与正文相关的图片
                if main_container:
                    for img in main_container.find_all('img', src=True):
                        img_url = urljoin(url, img['src'].strip())
                        if img_url not in images:
                            images.append(img_url)
                return {
                    'title': soup.title.string.strip() if soup.title else "无标题",
                    'text': text,
                    'html': str(soup),
                    'top_image': images[0] if images else None,
                    'images': images
                }
            except Exception as e2:
                logging.error(f"BeautifulSoup解析失败: {str(e2)}")
                raise

    def clean_text(self, text: str) -> str:
        """清理文本内容"""
        text = re.sub(r'\s+', ' ', text).strip()  # 移除多余空白
        text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)  # 移除特殊字符
        text = re.sub(r'([!?,.])\1+', r'\1', text)  # 移除重复标点
        return text

    def extract_main_content(self, soup: BeautifulSoup) -> Optional[BeautifulSoup]:
        """提取正文内容，优先选择文章、section、article-content等常用容器"""
        selectors = ['article', 'section', 'main', '.article-content', '.post-content', '#content']
        for sel in selectors:
            candidate = soup.select_one(sel)
            if candidate and len(candidate.get_text(strip=True)) > 500:
                return candidate
        # 后备方案：遍历多个包含较多文字的div或p标签
        paragraphs = soup.find_all(['p', 'div'])
        best_candidate = None
        max_text = ""
        for p in paragraphs:
            current_text = p.get_text(separator=" ", strip=True)
            if len(current_text) > len(max_text):
                max_text = current_text
                best_candidate = p
        return best_candidate if best_candidate else soup.body

    def analyze_structure(self, text: str) -> List[Dict]:
        """分析文章结构"""
        system_prompt = """你是一位专业的内容分析师，请按照以下要求分析文章结构：
1. 识别文章核心主题，主题的数量不超过20个
2. 识别每个核心主题所包含的内容要点，内容要点的数量不超过10个
3. 使用严格的JSON格式返回结果，只返回JSON，不要有任何额外文字

输出格式示例：
{
    "sections": [
        {
            "title": "主题标题",
            "summary": "主题摘要",
            "key_points": ["要点1", "要点2"]
        }
    ]
}"""
        
        user_prompt = f"""并返回严格JSON格式的结构化结果：
{text[:MAX_CONTENT_LENGTH]}"""
        
        try:
            response = self.call_deepseek_api(system_prompt, user_prompt)
            #print(response)
            if not response:
                return self.fallback_structure_analysis(text)
            
            # 尝试提取可能的JSON部分
            json_str = response.strip()
            if not (json_str.startswith('{') and json_str.endswith('}')):
                match = re.search(r'\{.*\}', json_str, re.DOTALL)
                if match:
                    json_str = match.group(0)
            
            data = json.loads(json_str)
            if not isinstance(data, dict) or 'sections' not in data:
                raise ValueError("无效的响应结构")
                
            valid_sections = []
            for section in data['sections']:
                if isinstance(section, dict) and 'title' in section:
                    key_points = section.get('key_points', [])
                    if not isinstance(key_points, list):
                        key_points = [str(key_points)] if key_points else []
                    valid_sections.append({
                        "title": str(section['title']),
                        "summary": str(section.get('summary', '')),
                        "key_points": key_points
                    })
                    
            if not valid_sections:
                raise ValueError("没有有效章节")
                
            return valid_sections
            
        except Exception as e:
            logging.error(f"结构解析失败: {str(e)}")
            logging.debug(f"原始响应: {response}")
            return self.fallback_structure_analysis(text)

    def fallback_structure_analysis(self, text: str) -> List[Dict]:
        """后备结构分析方案"""
        paragraphs = [p.strip() for p in text.split('\n') if len(p.strip()) > MIN_SECTION_LENGTH]
        sections = []
        for i, para in enumerate(paragraphs[:10]):  # 最多取10段
            num_points = min(3, len(para.split('。')))
            key_points = [f"要点 {j+1}" for j in range(num_points)]
            sections.append({
                "title": f"主题 {i+1}",
                "summary": para[:200] + "..." if len(para) > 200 else para,
                "key_points": key_points
            })
        return sections

    def create_ppt(self, url: str, output_path: str) -> bool:
        """主PPT生成函数"""
        try:
            logging.info(f"开始解析网页: {url}")
            article_data = self.parse_article(url)
            if not article_data or not article_data['text']:
                raise ValueError("无法获取有效文章内容")
                
            logging.info("分析文章结构...")
            structure = self.analyze_structure(article_data['text'])
            
            prs = Presentation()
            self.add_title_slide(prs, article_data['title'], url)
            
            for section in structure:
                self.add_content_slide(prs, section)
                
            if article_data['images']:
                self.add_image_slides(prs, article_data['images'], article_data['text'], url)
                
            prs.save(output_path)
            logging.info(f"PPT生成成功: {output_path}")
            return True
            
        except Exception as e:
            logging.error(f"PPT生成失败: {str(e)}")
            return False

    def add_title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        """标题幻灯片创建函数"""
        try:
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
            
            title_left = (prs.slide_width - Inches(8)) / 2
            title_box = slide.shapes.add_textbox(
                left=title_left,
                top=Inches(1),
                width=Inches(8),
                height=Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = title[:100]
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.font.name = '微软雅黑'
            
            
            subtitle_box = slide.shapes.add_textbox(
                left=title_left,
                top=Inches(2.5),
                width=Inches(8),
                height=Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.margin_top = Pt(50)
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle[:200]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.italic = True
            
        except Exception as e:
            logging.error(f"创建标题幻灯片失败: {str(e)}")
            raise RuntimeError(f"标题幻灯片创建错误: {str(e)}")

    def add_content_slide(self, prs: Presentation, section: Dict) -> None:
        """正文内容幻灯片创建函数"""
        try:
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            title_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(0.25),
                width=prs.slide_width - Inches(1),
                height=Inches(1))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = section.get('title', '无标题')[:100]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.name = '微软雅黑'
            
            content_box = slide.shapes.add_textbox(
                left=Inches(0.75),
                top=Inches(1.5),
                width=prs.slide_width - Inches(1.5),
                height=prs.slide_height - Inches(2))
            content_frame = content_box.text_frame
            content_frame.fit_text(font_family='微软雅黑', bold=True)
            
            if 'summary' in section and section['summary']:
                summary_para = content_frame.add_paragraph()
                summary_para.text = section['summary'][:300]
                summary_para.word_wrap = True
                summary_para.font.size = Pt(24)
                summary_para.space_after = Pt(12)
            
            for point in section.get('key_points', [])[:7]:
                point_para = content_frame.add_paragraph()
                point_para.text = f"• {point[:200]}"
                point_para.word_wrap = True
                point_para.level = 1
                point_para.font.size = Pt(24)
                point_para.space_after = Pt(6)
                point_para.font.name = '微软雅黑'
                
        except Exception as e:
            logging.error(f"创建内容幻灯片失败: {str(e)}")
            raise RuntimeError(f"内容幻灯片创建错误: {str(e)}")

    def add_image_slides(self, prs: Presentation, image_urls: List[str], text: str, base_url: str) -> None:
        """图片幻灯片添加函数，优先选择与文本相关性高的图片"""
        best_images = self.select_best_images(image_urls, text, base_url)
        
        for img_url, img_data, dimensions in best_images[:MAX_IMAGES_PER_PPT]:
            try:
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                title_box = slide.shapes.add_textbox(
                    left=Inches(0.5),
                    top=Inches(0.25),
                    width=prs.slide_width - Inches(1),
                    height=Inches(0.5))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "相关图示"
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                
                img_width, img_height = dimensions
                aspect_ratio = img_width / img_height
                
                max_width = Inches(7)
                max_height = Inches(5)
                
                if aspect_ratio > (max_width / max_height):
                    width = min(max_width, Inches(img_width / 96))  # 假设96dpi
                    height = width / aspect_ratio
                else:
                    height = min(max_height, Inches(img_height / 96))
                    width = height * aspect_ratio
                    
                left = (prs.slide_width - width) / 2
                top = (prs.slide_height - height) / 2 + Inches(0.5)  # 为标题留空间
                
                slide.shapes.add_picture(img_data, left, top, width, height)
                
            except Exception as e:
                logging.error(f"添加图片幻灯片失败: {str(e)}")
                continue

    def select_best_images(self, image_urls: List[str], text: str, base_url: str) -> List[Tuple[str, BytesIO, Tuple[int, int]]]:
        """选择最相关的图片，通过OCR提取图片中的文字，与全文进行TF-IDF相似度匹配"""
        scored_images = []
        downloaded_images = []
        for img_url in image_urls[:20]:  # 最多处理20张图片
            result = self.download_image(img_url, base_url)
            if result:
                img_data, dimensions = result
                downloaded_images.append((img_url, img_data, dimensions))
                
        if not downloaded_images:
            return []
            
        if len(downloaded_images) <= MAX_IMAGES_PER_PPT:
            return downloaded_images
            
        try:
            vectorizer = TfidfVectorizer(max_features=100, stop_words='english')
            text_tfidf = vectorizer.fit_transform([text])
            
            for img_url, img_data, dimensions in downloaded_images:
                try:
                    with Image.open(img_data) as img:
                        ocr_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    img_tfidf = vectorizer.transform([ocr_text])
                    similarity = cosine_similarity(text_tfidf, img_tfidf)[0][0]
                    scored_images.append((img_url, img_data, dimensions, similarity))
                    img_data.seek(0)
                except Exception as e:
                    logging.warning(f"图片处理失败 {img_url}: {str(e)}")
                    continue
                    
            scored_images.sort(key=lambda x: x[3], reverse=True)
            return [(url, data, dim) for url, data, dim, _ in scored_images]
            
        except Exception as e:
            logging.error(f"图片评分失败: {str(e)}")
            return downloaded_images[:MAX_IMAGES_PER_PPT]

    def add_single_image_slide(self, prs: Presentation, img_url: str, 
                             img_data: BytesIO, dimensions: Tuple[int, int]) -> None:
        """添加单张图片幻灯片"""
        try:
            slide_layout = prs.slide_layouts[5]  # 仅标题幻灯片
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = "相关图示"
            title_shape.font.size = Pt(28)
            
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width, img_height = dimensions
            aspect_ratio = img_width / img_height
            
            max_width = Inches(7.5)
            max_height = Inches(5.5)
            
            if aspect_ratio > (max_width / max_height):
                width = min(max_width, Inches(img_width / 96))
                height = width / aspect_ratio
            else:
                height = min(max_height, Inches(img_height / 96))
                width = height * aspect_ratio
                
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2 + Inches(0.5)
            
            slide.shapes.add_picture(img_data, left, top, width, height)
            
            notes_slide = slide.notes_slide
            if notes_slide:
                notes_slide.notes_text_frame.text = f"图片来源: {img_url}"
                
        except Exception as e:
            logging.error(f"添加图片幻灯片失败: {str(e)}")

if __name__ == "__main__":
    try:
        url = sys.argv[1]
        if not url:
            url = input("请输入网页URL: ").strip()
        if not url.startswith(('http://', 'https://')):
            url = 'http://' + url
            
        #output_file = input("输出文件名 (默认: output.pptx): ").strip() or "output.pptx"
        output_file = "output.pptx"
        
        generator = PPTGenerator()
        success = generator.create_ppt(url, output_file)
        
        if success:
            print(f"\nPPT生成成功: {output_file}")
            print(f"文件大小: {os.path.getsize(output_file)/1024:.1f} KB")
        else:
            print("\nPPT生成失败，请查看日志文件获取详细信息")
            
    except KeyboardInterrupt:
        print("\n操作已取消")
    except Exception as e:
        print(f"\n发生错误: {str(e)}")
